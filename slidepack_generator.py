"""
Slidepack Generator
Author(s):
Pri Balachandran

Module to simplify the creation of slidepacks.

Please note, these require your Powerpoint presentation to have slide layouts defined in the Master Slide view.
Slide layouts define placeholders for these modules to place content into.

Many modules you will not need to interact with directly, these are named with a leading underscore.

"""

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.slide import Slide


def get_all_slide_layouts(prs: Presentation):
    """
    Returns the slide layout name, mapped to slide number

    A 'slide layout' is effectively a template slide. When designing the whole slidepack template
    in Powerpoint, each slide layout (i.e. slide template) is given a name and slide number.
    This method allows you to retrieve the name and position of all slide layouts, without
    needing to open Powerpoint.Use this method when you need a reminder of the names of the
    slide layouts in your template deck, so you can pick the relevant one for add_slide.

    Args:
        prs: The presentation being used as a template

    Returns:
        dict: Slide layout name mapped to slide number
    """
    return {l.name: i for i, l in enumerate(prs.slide_layouts)}


def get_slide_placeholders(slide: Slide):
    """
    Provides a dictionary breakdown of the placeholders available on a given slide.

    This slidepack generator works by adding content into pre-defined placeholders. A placeholder
    is effectively the empty content spaces you see when clicking 'New Slide'. This method allows
    you to see what placeholders are on a slide, so you don't need to use Powerpoint.
    See the examples section for the output strucure

    Args:
        slide: Slide in pptx from which to find placeholders.

    Returns:
        dict: Placeholders indices, grouped into title, body, picture and other.


    Examples:
        Accessing slide from Presentation
        prs = Presentation(...)
        slide = prs.slides.add_slide(...)

        Example output, for a slide layout with a title and two pictures
        {'title': 0, 'body':[], 'picture': [7,8], 'other':{}}


    Notes:
        Valid placeholder types
        https://python-pptx.readthedocs.io/en/latest/api/enum/PpPlaceholderType.html

        Enums
        https://python-pptx.readthedocs.io/en/latest/dev/analysis/enumerations.html

        Shape ordering
        https://python-pptx.readthedocs.io/en/latest/api/shapes.html
    """

    slide_placeholders = {"title": "", "body": [], "picture": [], "other": {}}

    # Loop through shapes
    for shape in slide.shapes:
        if shape.is_placeholder:
            phf = shape.placeholder_format

            if phf.type == PP_PLACEHOLDER_TYPE.TITLE:
                slide_placeholders["title"] = int(phf.idx)
            elif phf.type == PP_PLACEHOLDER_TYPE.BODY:
                slide_placeholders["body"].append(int(phf.idx))
            elif phf.type == PP_PLACEHOLDER_TYPE.PICTURE:
                slide_placeholders["picture"].append(int(phf.idx))
            else:
                slide_placeholders["other"][phf.idx] = str(phf.type)

    # Shapes are ordered in reverse of the UI (i.e. starting from the bottom)
    # reverse it here so it goes from top-to-bottom
    slide_placeholders["body"].reverse()
    slide_placeholders["picture"].reverse()

    return slide_placeholders


def _get_slide_layout_idx(prs: Presentation, layout_name: str):
    """
    Returns the index of a layout in a presentation. Designed for use in other methods in this module,
    and may not be useful directly.

    Slide layouts are referenced by their index in most python-pptx methods. This method
    makes it easier for other methods to access that index, and ensure that anyone using the
    front-facing modules like add_slide, can use the layout_name instead of index.

    Args:
        prs: The presentation being used as a template
        layout_name: The name of the slide layout. This can be found in the Powerpoint Master Slide
            view or by using get_all_slide_layouts()

    Returns:
        int: The index of slide_layout
    """

    map_idx_slidelayout = get_all_slide_layouts(prs)

    return map_idx_slidelayout[layout_name]


def _calc_max_image_height_within_placeholder(
    img_size: (int, int), placeholder_size: (int, int)
):
    """
    Returns the maximal width and height for an image to fit into a placeholder, whilst maintaining
    the aspect ratio. Designed for use in other methods in this module, and may not be useful
    directly.

    This calculation provies the largest size of the image so neither height nor width exceeds that
    of the placeholder.

    Args:
        img_size: (width, height) Size of the image in px
        placeholder_size: (width, height) Size of the placeholder in px

    Returns:
        (int, int): (width, height) Size of the image where
            width <= placeholder width
            height <= placeholder height

    Examples:
        An image (10,40) to fit into a placeholder (40,80), returns (20, 80)
        An image (10,40) to fit into a placeholder (40,400). returns (40, 160)

    """
    img_width = img_size[0]
    img_height = img_size[1]

    placeholder_width = placeholder_size[0]
    placeholder_height = placeholder_size[1]

    width_scale_factor = placeholder_width / img_width
    height_scale_factor = placeholder_height / img_height

    # Get the actual scale factor - which is the smallest of these
    scale_factor = min(width_scale_factor, height_scale_factor)
    output_width = int(img_width * scale_factor)
    output_height = int(img_height * scale_factor)

    return (output_width, output_height)


def _remove_placeholder_from_slide(slide: Slide, placeholder_idx: int):
    """Updates slide object by removing a placeholder"""
    shape = slide.placeholders[placeholder_idx]
    shape = shape._element
    shape.getparent().remove(shape)


def _add_picture_within_placeholder(slide: Slide, placeholder_idx: int, picture: str):
    """
    Updates slide to have a picture in the position of a defined placeholder, without cropping.
    Designed for use in other methods in this module, and may not be useful directly.

    This function over comes the limitation of the python-pptx insert_picture method, by
    adding a picture without cropping, scaling it so that it is no larger than the placeholder.
    No returns. Changes are made to slide.

    Args:
        slide: Slide in pptx from which to find the picture placeholder.
        placeholder_idx: Index of the picture placeholder in slide.
        picture: Picture filepath as string.

    Notes:
        The limitation of the pptx insert_picture method, which this method overcomes,
        is outlined here- https://python-pptx.readthedocs.io/en/latest/user/
                        placeholders-using.html#pictureplaceholder-insert-picture
    """
    placeholder_left = slide.placeholders[placeholder_idx].left
    placeholder_top = slide.placeholders[placeholder_idx].top

    placeholder_width = slide.placeholders[placeholder_idx].width
    placeholder_height = slide.placeholders[placeholder_idx].height
    placeholder_size = (placeholder_width, placeholder_height)

    img_size = Image.open(picture).size

    # Calculate the size of the picture to be placed on slide
    output_size = _calc_max_image_height_within_placeholder(img_size, placeholder_size)

    slide.shapes.add_picture(
        picture,
        left=placeholder_left,
        top=placeholder_top,
        width=output_size[0],
        height=output_size[1],
    )

    _remove_placeholder_from_slide(slide, placeholder_idx)


def add_slide(
    prs: Presentation,
    layout_name: str,
    title="",
    bodies=[],
    pictures=[],
    picture_scale_method="fill_placeholder",
):
    """
    Adds a slide to a presentation, using a pre-defined layout.

    Design the layout using Powerpoint, ensuring that it has only Title, Body
    and Picture placeholders which this will populate. Details on how to do this
    are in the Notes section below. The body and pictures are put in order of their
    indices, which are generally in the order they were created.


    Args:
        prs: The presentation being used as a template (including any front pages to retain).
        layout_name: The layout to use in creating a slide. If you are unsure about what
            the names of the layouts are, use get_all_slide_layouts() or see the notes below.
        title: Title for the page. Leave as a blank string if the page has no title placeholder
        bodies: List of strings, one per text placeholder, in the order defined in Powerpoint Selection Pane
            (see Notes).
        pictures: List of local filepaths of PNGs. One for each picture placeholder, in the order defined in Powerpoint
        Selection Pane (see Notes).
        picture_scale_method: {'fill_placeholder', 'within_placeholder'}
            In all three scaling methods, the aspect ratio is maintained. The differences
            are around cropping and defining the end of the image.
            - 'fill_placeholder' will use the default Powerpoint logic and enlarge the image
            until the entire placeholder is filled. Some of the image may be cropped.
            - 'within_placeholder' will resize the image to be as large as possible, but
            staying entirely within the placeholder

    Returns:
        slide_idx: Integer representing the slide index. Retain for use in a contents page or similar

    Notes:
        This method requires the following:
        - A non-empty Title string is provided ONLY if there is a title placeholder.
        - Bodies must be empty if there are no placeholders for bodies
        - Bodies must contain no more strings than there are placeholders for bodies.
        - Pictures must be empty if there are no placeholders for pictures
        - Pictures must contain no more strings than there are placeholders for pictures.


    Example:
        pres = Presentation('filename.pptx')
        _slide2_idx = add_slide(...)
        _slide3_idx = add_slide(...)
        ...
        pres.save('output_file.pptx')

    """

    assert (
        type(title) == str
    ), "Ensure title is string (empty string if no Title Placeholder exists)"

    # Get slide
    layout_idx = _get_slide_layout_idx(prs, layout_name)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    # Get all place holders on slide
    placeholders = get_slide_placeholders(slide)

    # CHECK INPUTS
    # TITLE
    # Ensure no title is provided iff there is no placeholder
    if placeholders["title"] == "":
        assert title == "", "Title provided but no title placeholder exists"

    # BODY
    # Ensure there are no more input bodies than there are body placeholders
    num_body_placeholders = len(placeholders["body"])
    num_bodies = len(bodies)
    assert (
        num_bodies <= num_body_placeholders
    ), f"Body: {num_bodies} strings provided, {num_body_placeholders} placeholders exist."

    # PICTURE
    # Ensure there are no more input pictures than there are picture placeholders
    num_pic_placeholders = len(placeholders["picture"])
    num_pictures = len(pictures)
    assert (
        num_pictures <= num_pic_placeholders
    ), f"Pictures: {num_pictures} provided, but {num_pic_placeholders} placeholders exists."

    # UPDATE SLIDE

    # Update Title
    if placeholders["title"] != "":
        if title != "":
            slide_title = slide.shapes.title
            slide_title.text = title
        else:
            _remove_placeholder_from_slide(slide, placeholders["title"])

    # Update Body
    # Loop through placeholders
    for i, placeholder_idx in enumerate(placeholders["body"]):
        # If a string exists in body, use it to update the placeholder text
        if i < num_bodies:
            body_placeholder = slide.placeholders[placeholder_idx]
            body_placeholder.text = bodies[i]

        # else remove the placeholder
        else:
            _remove_placeholder_from_slide(slide, placeholder_idx)

    # Update Picture
    # Loop through placeholders
    for i, placeholder_idx in enumerate(placeholders["picture"]):
        # If a picture filepath exists, use it to insert a picture into the placeholder
        if i < num_pictures:

            picture = pictures[i]

            if picture_scale_method == "fill_placeholder":
                placeholder = slide.placeholders[
                    placeholder_idx
                ]  # idx key, not position
                placeholder.insert_picture(picture)

            elif picture_scale_method == "within_placeholder":
                _add_picture_within_placeholder(slide, placeholder_idx, picture)

        # Else remove the placeholder
        else:
            _remove_placeholder_from_slide(slide, placeholder_idx)

    return slide.slide_id
