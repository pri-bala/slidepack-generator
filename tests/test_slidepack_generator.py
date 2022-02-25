"""
Slidepack Generator Tests
Pri Balachandran

Requires slidepack_gen_test_layouts_20210126.pptx, which contains some 
slide layouts with known placeholders and structures to test against.
"""

import sys
import os
import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

import slidepack_generator as sg

sample_img_fname = "tests/support/slidepack_gen_test_image.png"


@pytest.fixture
# Get Presentation
def get_prs():
    print("setup")
    prs = Presentation("tests/support/slidepack_gen_test_layouts_20210126.pptx")
    yield prs
    print("teardown")


@pytest.fixture
# Get specific slide which contanins no  placeholders
def get_no_placeholders_slide(get_prs):
    # Slide layout index of slide with no placeholders
    layout_no_placeholders_idx = 3
    get_prs.add_slide(prs.slide_layouts[layout_no_placeholders_idx])


### TESTS ###

## METHOD get_all_slide_layouts ##
# Must have static slide_layout index and names to test against
# otherwise other tests will fail...
def test_get_all_slide_layouts_contains_known_names(get_prs):
    dct_layout_names = sg.get_all_slide_layouts(get_prs)
    assert ("no_placeholders", 0) in dct_layout_names.items()
    assert ("title_only", 1) in dct_layout_names.items()
    assert ("title_2txt", 2) in dct_layout_names.items()
    assert ("title_2pic_3txt", 3) in dct_layout_names.items()
    assert ("other_placeholders", 4) in dct_layout_names.items()


## METHOD _get_slide_layout_idx ##

# Must return integer
def test_get_slide_layout_idx_returns_int(get_prs):
    slide_layout_name = "no_placeholders"
    slide_layout_idx = sg._get_slide_layout_idx(get_prs, slide_layout_name)
    assert type(slide_layout_idx) == int


# METHOD get_slide_placeholders ##

# Must have same structure when empty
def test_get_slide_placeholders(get_prs):
    idx_no_placeholders = 0  # can be identified using sg._get_slide_layout_idx()
    slide = get_prs.slides.add_slide(get_prs.slide_layouts[idx_no_placeholders])
    placeholders = sg.get_slide_placeholders(slide)

    expected_structure = {"title": "", "body": [], "picture": [], "other": {}}
    assert placeholders == expected_structure


# Must have same sub-structure when some placeholders exist
def test_get_slide_placeholders_title_2pic_3txt(get_prs):
    idx_title_2pic_3txt = 3  # can be identified using sg._get_slide_layout_idx()
    slide = get_prs.slides.add_slide(get_prs.slide_layouts[idx_title_2pic_3txt])
    placeholders = sg.get_slide_placeholders(slide)

    assert placeholders["title"] != ""
    assert len(placeholders["picture"]) == 2
    assert len(placeholders["body"]) == 3
    assert len(placeholders["other"]) == 0


# Must have same sub-structure when other content exists
def test_get_slide_placeholders_other_placeholders(get_prs):
    idx_other_placeholders = 4  # can be identified using sg._get_slide_layout_idx()
    slide = get_prs.slides.add_slide(get_prs.slide_layouts[idx_other_placeholders])
    placeholders = sg.get_slide_placeholders(slide)

    assert placeholders["title"] != ""
    assert len(placeholders["picture"]) == 0
    assert len(placeholders["body"]) == 0
    assert len(placeholders["other"]) == 1


## METHOD add_picture_within_placeholder
# Ensure placeholder idx is removed and picture idx is added
def test_add_picure_within_placeholder(get_prs):
    idx_no_placeholders = 3  # can be identified using sg._get_slide_layout_idx()
    slide = get_prs.slides.add_slide(get_prs.slide_layouts[idx_no_placeholders])

    # Get placeholder indexes before method
    pre_placeholders = sg.get_slide_placeholders(slide)
    pre_idxs = pre_placeholders["picture"]

    remove_idx = pre_idxs[0]

    # Method being tested
    sg._add_picture_within_placeholder(slide, remove_idx, sample_img_fname)

    # Get placeholder indexes before method
    post_placeholders = sg.get_slide_placeholders(slide)
    post_idxs = post_placeholders["picture"]
    new_idx = post_idxs[-1]

    assert remove_idx not in post_placeholders
    assert new_idx not in pre_placeholders


def test_calc_max_image_height_within_placeholder():
    img_size = (1000, 500)  # width, height in px
    image = Image.new("RGB", img_size)

    # width x 1.2
    placeholder_size = (1200, 2000)
    output_size = sg._calc_max_image_height_within_placeholder(
        image.size, placeholder_size
    )
    assert output_size == (1200, 600)

    # height x 2
    placeholder_size = (3000, 1000)
    output_size = sg._calc_max_image_height_within_placeholder(
        image.size, placeholder_size
    )
    assert output_size == (2000, 1000)

    # both
    placeholder_size = (4000, 2000)
    output_size = sg._calc_max_image_height_within_placeholder(
        image.size, placeholder_size
    )
    assert output_size == (4000, 2000)


## FUNCTION add_slide ##

# TITLES
# No placeholder, must allow empty title string
def test_title_not_accepted_not_provided(get_prs):
    slide_idx = sg.add_slide(prs=get_prs, layout_name="no_placeholders", title="")
    assert type(slide_idx) == int


# No placeholder, must throw error if title provided
def test_title_not_accepted_but_provided(get_prs):
    with pytest.raises(AssertionError):
        sg.add_slide(
            prs=get_prs, layout_name="no_placeholders", title="Some Title Text"
        )


# Placeholder exists, must allow nonempty title string to be passed
def test_title_accepted_and_provided(get_prs):
    slide_idx = sg.add_slide(
        prs=get_prs, layout_name="title_only", title="Title Text"
    )  # Pass desired
    assert type(slide_idx) == int


## BODY / TEXT
# No placeholder, must allow empty body to be provided
def test_body_not_accepted_not_provided(get_prs):
    slide_idx = sg.add_slide(
        prs=get_prs, layout_name="title_only", title="Some Title", bodies=[]
    )
    assert type(slide_idx) == int


# No placeholder, must throw error if body provided
def test_body_not_accepted_but_provided(get_prs):
    with pytest.raises(AssertionError):
        sg.add_slide(
            prs=get_prs,
            layout_name="title_only",
            title="Some Title",
            bodies=["some text1"],
        )


# Placeholder exists, must allow some body to be provided
def test_body_accepted_and_provided(get_prs):
    slide_idx = sg.add_slide(
        prs=get_prs,
        layout_name="title_2txt",
        title="Some Title",
        bodies=["some text1", "some text2"],
    )
    assert type(slide_idx) == int


# Arg given as string must fail
def test_body_accepted_given_as_string(get_prs):
    with pytest.raises(AssertionError):
        sg.add_slide(
            prs=get_prs,
            layout_name="title_2txt",
            title="Some Title",
            bodies="some text1",
        )


# Placeholder(s) exist, but fail as too many bodies provided
def test_body_accepted_too_many_provided(get_prs):
    with pytest.raises(AssertionError):
        sg.add_slide(
            prs=get_prs,
            layout_name="title_2txt",
            title="Some Title",
            bodies=["some text1", "some text2", "some text3"],
        )


## PICTURES
# No placeholder, must allow empty picture to be provided
def test_picture_not_accepted_not_provided(get_prs):
    slide_idx = sg.add_slide(
        prs=get_prs, layout_name="title_only", title="Some Title", pictures=[]
    )
    assert type(slide_idx) == int


# No placeholder, must throw error if picture provided
def test_body_not_accepted_but_provided(get_prs):
    with pytest.raises(AssertionError):
        sg.add_slide(
            prs=get_prs,
            layout_name="title_only",
            title="Some Title",
            pictures=[sample_img_fname],
        )


# Placeholder exists, must allow some picture to be provided
def test_picture_accepted_and_provided(get_prs):
    slide_idx = sg.add_slide(
        prs=get_prs,
        layout_name="title_2pic_3txt",
        title="Some Title",
        pictures=[sample_img_fname],
    )
    assert type(slide_idx) == int


# Arg given as string must fail
def test_picture_accepted_given_as_string(get_prs):
    with pytest.raises(AssertionError):
        sg.add_slide(
            prs=get_prs,
            layout_name="title_2pic_3txt",
            title="Some Title",
            pictures=sample_img_fname,
        )


# Placeholder(s) exist, but fail as too many pictures provided
def test_picture_accepted_too_many_provided(get_prs):
    with pytest.raises(AssertionError):
        sg.add_slide(
            prs=get_prs,
            layout_name="title_2pic_3txt",
            title="Some Title",
            pictures=[sample_img_fname for i in range(3)],
        )


def test_add_slide_result(get_prs):
    prs = get_prs

    new_slide_idx = sg.add_slide(
        prs=prs,
        layout_name="title_2pic_3txt",
        title="Some Title",
        bodies=["text1", "text2"],
        pictures=[sample_img_fname],
    )
    # Get newly created slide
    slide = prs.slides.get(new_slide_idx)

    placeholders = sg.get_slide_placeholders(slide)

    assert placeholders["title"] != ""
    assert len(placeholders["body"]) == 2
    assert len(placeholders["picture"]) == 1


def test_add_slide_within_placeholder_result(get_prs):
    prs = get_prs

    new_slide_idx = sg.add_slide(
        prs=prs,
        layout_name="title_2pic_3txt",
        title="Some Title",
        bodies=["text1", "text2"],
        pictures=[sample_img_fname],
        picture_scale_method="within_placeholder",
    )
    # Get newly created slide
    slide = prs.slides.get(new_slide_idx)

    placeholders = sg.get_slide_placeholders(slide)

    assert placeholders["title"] != ""
    assert len(placeholders["body"]) == 2

    # As using 'within placeholder' replaces the placeholder with a picture...
    num_pics = sum(
        [shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes]
    )
    assert num_pics == 1
